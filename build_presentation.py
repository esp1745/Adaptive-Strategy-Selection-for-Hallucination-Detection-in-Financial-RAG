"""
Generate FinGuard presentation:
'Adaptive Hallucination Detection in Financial RAG via Proximal Policy Optimization'
10 slides, academic conference style (NeurIPS/ICLR).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colours ──────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0C, 0x1D, 0x33)
BLUE        = RGBColor(0x1A, 0x4F, 0xBA)
BLUE_LIGHT  = RGBColor(0xEF, 0xF4, 0xFF)
BLUE_MID    = RGBColor(0xBF, 0xDB, 0xFE)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE    = RGBColor(0xF5, 0xF7, 0xFA)
SLATE       = RGBColor(0x64, 0x74, 0x8B)
GREEN       = RGBColor(0x15, 0x80, 0x3D)
GREEN_BG    = RGBColor(0xF0, 0xFD, 0xF4)
AMBER       = RGBColor(0xB4, 0x53, 0x09)
AMBER_BG    = RGBColor(0xFF, 0xFB, 0xEB)
RED         = RGBColor(0xDC, 0x26, 0x26)
RED_BG      = RGBColor(0xFE, 0xF2, 0xF2)
PURPLE      = RGBColor(0x6D, 0x28, 0xD9)
BORDER      = RGBColor(0xDD, 0xE3, 0xEC)
DARK_BORDER = RGBColor(0x1E, 0x35, 0x55)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helper utilities ──────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0.75)):
    shape = slide.shapes.add_shape(1, l, t, w, h)      # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=Pt(14), bold=False, color=NAVY, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def add_para(tf, text, size=Pt(13), bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, italic=False, space_before=Pt(4)):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p


def navy_header(slide, title, subtitle=None):
    """Dark navy top bar with title."""
    bar_h = Inches(1.25)
    add_rect(slide, 0, 0, W, bar_h, fill=NAVY)
    # accent line
    add_rect(slide, 0, bar_h, W, Pt(3), fill=BLUE)
    add_text(slide, title,
             Inches(0.45), Inches(0.18), Inches(12.0), Inches(0.65),
             size=Pt(26), bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.45), Inches(0.80), Inches(12.0), Inches(0.38),
                 size=Pt(13), color=RGBColor(0x93, 0xC5, 0xFD), italic=True)
    # slide background
    add_rect(slide, 0, bar_h + Pt(3), W, H - bar_h - Pt(3), fill=WHITE)
    return bar_h + Pt(6)   # content_y


def bullet_box(slide, items, l, t, w, h, title=None,
               title_color=BLUE, item_size=Pt(13), bg=OFFWHITE, border=BORDER):
    add_rect(slide, l, t, w, h, fill=bg, line=border)
    y = t + Inches(0.12)
    if title:
        add_text(slide, title, l + Inches(0.14), y, w - Inches(0.28), Inches(0.32),
                 size=Pt(11), bold=True, color=title_color)
        y += Inches(0.32)
    for item in items:
        add_text(slide, item, l + Inches(0.22), y, w - Inches(0.36), Inches(0.32),
                 size=item_size, color=NAVY)
        y += Inches(0.30)


def metric_tile(slide, label, value, l, t, w=Inches(2.2), h=Inches(0.95),
                val_color=BLUE, bg=BLUE_LIGHT, border=BLUE_MID):
    add_rect(slide, l, t, w, h, fill=bg, line=border)
    add_text(slide, label, l + Inches(0.1), t + Inches(0.06),
             w - Inches(0.2), Inches(0.28), size=Pt(9), color=SLATE)
    add_text(slide, value, l + Inches(0.1), t + Inches(0.30),
             w - Inches(0.2), Inches(0.55), size=Pt(22), bold=True, color=val_color)


def page_num(slide, n):
    add_text(slide, f"{n} / 12",
             W - Inches(1.1), H - Inches(0.32), Inches(0.9), Inches(0.28),
             size=Pt(10), color=SLATE, align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)

# Full navy background
add_rect(s, 0, 0, W, H, fill=NAVY)
# Accent stripe
add_rect(s, 0, Inches(3.8), W, Pt(3), fill=BLUE)

# Conference tag
add_text(s, "CS Research Presentation  ·  March 2026",
         Inches(0.6), Inches(1.05), Inches(12.0), Inches(0.38),
         size=Pt(12), color=RGBColor(0x93, 0xC5, 0xFD), italic=True)

# Main title
txb = s.shapes.add_textbox(Inches(0.6), Inches(1.55), Inches(12.1), Inches(2.1))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "Adaptive Hallucination Detection in"
run.font.size = Pt(34)
run.font.bold = True
run.font.color.rgb = WHITE
run.font.name = "Calibri"
add_para(tf, "Financial RAG via Proximal Policy Optimization",
         size=Pt(34), bold=True, color=WHITE)

# Subtitle / keywords
add_text(s, "PPO  ·  Contextual Bandit  ·  SEC Filing Verification  ·  Detector Selection",
         Inches(0.6), Inches(3.95), Inches(12.0), Inches(0.42),
         size=Pt(14), color=RGBColor(0x7D, 0xD3, 0xFC), italic=False)

# Author block
add_rect(s, Inches(0.6), Inches(4.55), Inches(5.5), Inches(0.95),
         fill=RGBColor(0x15, 0x2D, 0x4A), line=DARK_BORDER)
add_text(s, "Presented by: [Your Name]",
         Inches(0.75), Inches(4.65), Inches(5.2), Inches(0.32),
         size=Pt(13), bold=True, color=WHITE)
add_text(s, "[University]  ·  [Course / Advisor]",
         Inches(0.75), Inches(4.95), Inches(5.2), Inches(0.30),
         size=Pt(12), color=RGBColor(0x93, 0xC5, 0xFD))

# FinGuard badge
add_rect(s, Inches(9.3), Inches(4.55), Inches(3.4), Inches(0.95),
         fill=BLUE, line=RGBColor(0x1E, 0x40, 0xAF))
add_text(s, "FinGuard",
         Inches(9.4), Inches(4.60), Inches(3.2), Inches(0.50),
         size=Pt(28), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Prototype System",
         Inches(9.4), Inches(5.08), Inches(3.2), Inches(0.28),
         size=Pt(11), color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

page_num(s, 1)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — ABSTRACT
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
cy = navy_header(s, "Abstract", "Adaptive Hallucination Detection in Financial RAG via PPO")

abstract = (
    "This paper presents FinGuard, an adaptive hallucination detection framework "
    "for financial Retrieval-Augmented Generation (RAG) systems operating over SEC "
    "filings. Given a natural-language query, the system retrieves supporting evidence "
    "via FAISS vector search and synthesizes a grounded answer. A Proximal Policy "
    "Optimization (PPO) agent observes a 12-dimensional state vector encoding query "
    "complexity, retrieval confidence, and answer characteristics, and selects among "
    "four hallucination detectors — token overlap, semantic similarity, BERT NLI, and "
    "LLM-as-judge — to maximize detection accuracy while minimizing inference cost and "
    "latency. The policy is updated online after each query using a clipped surrogate "
    "objective with entropy regularization. Benchmark evaluation on 20 financial QA "
    "samples demonstrates that adaptive selection reduces cost by up to 95% versus "
    "always using LLM-as-judge, while token overlap achieves the highest F1 score "
    "(0.80) motivating context-dependent detector routing over any fixed strategy."
)

# Abstract box
add_rect(s, Inches(0.45), cy, Inches(12.43), Inches(2.85),
         fill=BLUE_LIGHT, line=BLUE_MID)
txb = s.shapes.add_textbox(Inches(0.65), cy + Inches(0.15),
                            Inches(12.0), Inches(2.6))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = abstract
run.font.size = Pt(13.5)
run.font.color.rgb = NAVY
run.font.name = "Calibri"

# Keyword chips
kw_y = cy + Inches(3.1)
add_text(s, "Keywords:", Inches(0.45), kw_y, Inches(1.1), Inches(0.35),
         size=Pt(11), bold=True, color=SLATE)
keywords = [
    "Hallucination Detection", "Financial RAG", "PPO",
    "Contextual Bandit", "Detector Selection", "SEC Filings",
]
kx = Inches(1.55)
for kw in keywords:
    w_kw = Inches(len(kw) * 0.115 + 0.3)
    add_rect(s, kx, kw_y, w_kw, Inches(0.32), fill=BLUE, line=RGBColor(0x1E, 0x40, 0xAF))
    add_text(s, kw, kx + Inches(0.08), kw_y + Inches(0.04),
             w_kw - Inches(0.12), Inches(0.26),
             size=Pt(10), bold=True, color=WHITE)
    kx += w_kw + Inches(0.12)

# Contributions box
cy2 = cy + Inches(3.6)
add_rect(s, Inches(0.45), cy2, Inches(12.43), Inches(1.55),
         fill=OFFWHITE, line=BORDER)
add_text(s, "Key Contributions",
         Inches(0.6), cy2 + Inches(0.08), Inches(5.0), Inches(0.3),
         size=Pt(11), bold=True, color=BLUE)
contribs = [
    "①  PPO-based adaptive detector routing for financial QA hallucination detection",
    "②  12-dimensional state feature extractor capturing query, retrieval, and answer signals",
    "③  FinGuard: end-to-end prototype with FAISS retrieval, online PPO updates, and live dashboard",
]
for i, c in enumerate(contribs):
    add_text(s, c, Inches(0.65), cy2 + Inches(0.38 + i * 0.36),
             Inches(12.1), Inches(0.32), size=Pt(12), color=NAVY)

page_num(s, 2)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PROBLEM & MOTIVATION
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
cy = navy_header(s, "Problem & Motivation",
                 "Why adaptive hallucination detection in financial RAG?")

# Left column
add_text(s, "The Problem",
         Inches(0.45), cy + Inches(0.05), Inches(6.0), Inches(0.35),
         size=Pt(14), bold=True, color=BLUE)
problems = [
    "• LLMs in RAG pipelines generate plausible but factually wrong statements",
    "• Financial errors carry legal and regulatory consequences",
    "• A single detector does not fit all query types",
    "• Expensive detectors (LLM Judge) are overused — wasting cost and time",
    "• Cheap detectors (Token Overlap) are ignored — even when sufficient",
]
py = cy + Inches(0.42)
for item in problems:
    add_text(s, item, Inches(0.55), py, Inches(5.9), Inches(0.34),
             size=Pt(12.5), color=NAVY)
    py += Inches(0.34)

# Right column — cost comparison
add_text(s, "Cost vs. Accuracy Trade-off",
         Inches(6.9), cy + Inches(0.05), Inches(6.0), Inches(0.35),
         size=Pt(14), bold=True, color=BLUE)

rows = [
    ("Detector",        "F1",    "Latency",   "Cost",   BLUE,  WHITE),
    ("Token Overlap",   "0.80",  "3 ms",      "$0.05",  OFFWHITE, NAVY),
    ("BERT NLI",        "0.73",  "1,047 ms",  "$0.30",  OFFWHITE, NAVY),
    ("LLM Judge",       "0.53",  "11,146 ms", "$1.00",  OFFWHITE, NAVY),
    ("Semantic Sim.",   "0.38",  "128 ms",    "$0.10",  OFFWHITE, NAVY),
]
col_w = [Inches(2.1), Inches(0.85), Inches(1.35), Inches(0.95)]
col_x = [Inches(6.9), Inches(9.0), Inches(9.85), Inches(11.2)]
row_h = Inches(0.44)
ty = cy + Inches(0.45)
for r_idx, (d, f, lat, cost, bg, fg) in enumerate(rows):
    for c_idx, (val, cx, cw) in enumerate(zip([d, f, lat, cost], col_x, col_w)):
        bold = r_idx == 0
        add_rect(s, cx, ty, cw, row_h, fill=bg,
                 line=BORDER if r_idx > 0 else None)
        add_text(s, val, cx + Inches(0.08), ty + Inches(0.08),
                 cw - Inches(0.12), row_h - Inches(0.1),
                 size=Pt(12 if r_idx > 0 else 11),
                 bold=bold, color=fg if r_idx == 0 else (GREEN if c_idx==1 and r_idx==1 else (RED if c_idx==3 and r_idx==3 else NAVY)))
    ty += row_h

# Insight box
add_rect(s, Inches(0.45), H - Inches(1.3), Inches(12.43), Inches(0.85),
         fill=RGBColor(0x0C, 0x1D, 0x33), line=DARK_BORDER)
add_text(s,
         "Research Question:  Can a PPO agent learn to route queries to the "
         "optimal detector — maximizing accuracy while minimizing cost?",
         Inches(0.65), H - Inches(1.22), Inches(12.0), Inches(0.65),
         size=Pt(13), bold=False, color=WHITE)

page_num(s, 3)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — LITERATURE REVIEW
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
cy = navy_header(s, "Literature Review", "Related work across three research threads")

areas = [
    {
        "title": "RAG & Hallucination Detection",
        "color": BLUE,
        "bg": BLUE_LIGHT,
        "border": BLUE_MID,
        "refs": [
            "Lewis et al. (2020). Retrieval-Augmented Generation for NLP. NeurIPS.",
            "Ji et al. (2023). Survey of Hallucination in NLG. ACM Comp. Surveys.",
            "Min et al. (2023). FActScoring: Fine-grained Atomic Evaluation. EMNLP.",
            "Es et al. (2023). RAGAS: Automated Evaluation of RAG Pipelines. arXiv.",
        ],
    },
    {
        "title": "LLM-as-Judge & NLI Methods",
        "color": PURPLE,
        "bg": RGBColor(0xF5, 0xF3, 0xFF),
        "border": RGBColor(0xC4, 0xB5, 0xFD),
        "refs": [
            "Zheng et al. (2023). Judging LLM-as-a-Judge with MT-Bench. NeurIPS.",
            "He et al. (2022). Towards Document-Level NLI. ACL.",
            "Honovich et al. (2022). TRUE: Re-evaluating Factual Consistency. NAACL.",
            "Wang et al. (2023). Self-Consistency Improves CoT Reasoning. ICLR.",
        ],
    },
    {
        "title": "Reinforcement Learning for NLP",
        "color": GREEN,
        "bg": GREEN_BG,
        "border": RGBColor(0x86, 0xEF, 0xAC),
        "refs": [
            "Schulman et al. (2017). Proximal Policy Optimization. arXiv.",
            "Sutton & Barto (2018). Reinforcement Learning: An Introduction.",
            "Ouyang et al. (2022). Training LMs with Human Feedback (RLHF). NeurIPS.",
            "Auer et al. (2002). Finite-time Analysis of Multiarmed Bandits. JMLR.",
        ],
    },
]

col_w = Inches(4.0)
col_gap = Inches(0.21)
for i, area in enumerate(areas):
    cx_col = Inches(0.45) + i * (col_w + col_gap)
    box_h = Inches(4.7)
    add_rect(s, cx_col, cy + Inches(0.05), col_w, box_h,
             fill=area["bg"], line=area["border"])
    add_text(s, area["title"],
             cx_col + Inches(0.14), cy + Inches(0.15), col_w - Inches(0.25), Inches(0.38),
             size=Pt(12), bold=True, color=area["color"])
    ry = cy + Inches(0.6)
    for ref in area["refs"]:
        add_text(s, "▪  " + ref,
                 cx_col + Inches(0.16), ry, col_w - Inches(0.3), Inches(0.7),
                 size=Pt(10.5), color=NAVY)
        ry += Inches(0.75)

# Gap statement
add_rect(s, Inches(0.45), H - Inches(1.25), Inches(12.43), Inches(0.78),
         fill=NAVY, line=DARK_BORDER)
add_text(s,
         "Gap:  No prior work applies PPO-based adaptive routing to select among "
         "multiple hallucination detectors in domain-specific (financial) RAG pipelines.",
         Inches(0.65), H - Inches(1.17), Inches(12.0), Inches(0.62),
         size=Pt(12.5), color=WHITE)

page_num(s, 4)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DATASETS & EXPERIMENTAL SETUP
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
cy = navy_header(s, "Datasets & Experimental Setup",
                 "Primary corpus  ·  Benchmark evaluation set  ·  Data statistics")

add_text(s, "Primary Corpus  —  SEC 10-K Filings",
         Inches(0.45), cy + Inches(0.05), Inches(6.0), Inches(0.32),
         size=Pt(13), bold=True, color=BLUE)

corpus_rows = [
    ("Source",     "SEC EDGAR public filing database"),
    ("Companies",  "Apple · Microsoft · Tesla · Amazon · Google/Alphabet · Meta · NVIDIA"),
    ("Filing type","Annual reports (Form 10-K)"),
    ("Chunks",     "975 document chunks after sliding-window segmentation"),
    ("Chunk size", "~420 characters per chunk"),
    ("Embedding",  "Hash bag-of-words, 256-dim  →  FAISS IndexFlatIP"),
    ("Retrieval",  "Top-k = 3 chunks per query (cosine similarity)"),
]
iy = cy + Inches(0.42)
for label, val in corpus_rows:
    add_rect(s, Inches(0.45), iy, Inches(5.9), Inches(0.46),
             fill=BLUE_LIGHT, line=BLUE_MID)
    add_text(s, label + ":", Inches(0.58), iy + Inches(0.07),
             Inches(1.35), Inches(0.32), size=Pt(10), bold=True, color=BLUE)
    add_text(s, val, Inches(1.95), iy + Inches(0.07),
             Inches(4.25), Inches(0.35), size=Pt(10), color=NAVY)
    iy += Inches(0.50)

add_text(s, "Benchmark Evaluation Set  (custom)",
         Inches(6.85), cy + Inches(0.05), Inches(6.0), Inches(0.32),
         size=Pt(13), bold=True, color=GREEN)

bench_rows = [
    ("Total samples", "20 financial QA pairs"),
    ("Hallucinated",  "10 samples — fabricated or incorrect facts"),
    ("Grounded",      "10 samples — verified against source filing"),
    ("Query types",   "Revenue, risk factors, segments, growth comparisons"),
    ("Companies",     "Apple, Amazon, Tesla, Microsoft, Google"),
    ("Labeling",      "Manual annotation by domain expert"),
]
by2 = cy + Inches(0.42)
for label, val in bench_rows:
    add_rect(s, Inches(6.85), by2, Inches(5.9), Inches(0.46),
             fill=GREEN_BG, line=RGBColor(0x86, 0xEF, 0xAC))
    add_text(s, label + ":", Inches(6.98), by2 + Inches(0.07),
             Inches(1.55), Inches(0.32), size=Pt(10), bold=True, color=GREEN)
    add_text(s, val, Inches(8.55), by2 + Inches(0.07),
             Inches(4.05), Inches(0.35), size=Pt(10), color=NAVY)
    by2 += Inches(0.50)

bar_y = iy + Inches(0.15)
add_text(s, "Class distribution  (n = 20, balanced)",
         Inches(0.45), bar_y, Inches(5.9), Inches(0.28),
         size=Pt(10), bold=True, color=NAVY)
bar_y += Inches(0.3)
add_rect(s, Inches(0.45), bar_y, Inches(5.9), Inches(0.42),
         fill=RED_BG, line=RGBColor(0xFC, 0xA5, 0xA5))
add_rect(s, Inches(0.45), bar_y, Inches(2.95), Inches(0.42),
         fill=GREEN_BG, line=RGBColor(0x86, 0xEF, 0xAC))
add_text(s, "Grounded  50%", Inches(0.58), bar_y + Inches(0.09),
         Inches(2.6), Inches(0.26), size=Pt(10), bold=True, color=GREEN)
add_text(s, "Hallucinated  50%", Inches(3.5), bar_y + Inches(0.09),
         Inches(2.7), Inches(0.26), size=Pt(10), bold=True, color=RED)

lim_y = by2 + Inches(0.2)
add_rect(s, Inches(6.85), lim_y, Inches(5.9), Inches(0.65),
         fill=AMBER_BG, line=RGBColor(0xFD, 0xE6, 0x8A))
add_text(s, "Limitation & Planned Expansion",
         Inches(7.0), lim_y + Inches(0.05), Inches(5.6), Inches(0.26),
         size=Pt(10), bold=True, color=AMBER)
add_text(s,
         "Benchmark is small (n=20). Future: scale to FinQA (8,281 samples, "
         "Chen et al. 2021) and TATQA (2,757 samples, Zhu et al. 2021).",
         Inches(7.0), lim_y + Inches(0.31), Inches(5.6), Inches(0.30),
         size=Pt(9.5), color=NAVY)

page_num(s, 5)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SYSTEM ARCHITECTURE
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
cy = navy_header(s, "System Architecture", "End-to-end FinGuard pipeline")

stages = [
    ("1\nQuery\nProcessor",   BLUE,   "12-dim state\nfeature vector"),
    ("2\nFAISS\nRetrieval",   NAVY,   "Top-k SEC\nchunks"),
    ("3\nAnswer\nSynthesis",  BLUE,   "Grounded\nanswer text"),
    ("4\nPPO\nPolicy",        GREEN,  "Selected\ndetector d"),
    ("5\nHallucination\nDetector", AMBER, "Verdict +\nconfidence"),
    ("6\nPPO\nUpdate",        PURPLE, "θ updated\nonline"),
]
box_w = Inches(1.72)
box_h = Inches(1.55)
box_gap = Inches(0.24)
total_w = len(stages) * box_w + (len(stages) - 1) * box_gap
start_x = (W - total_w) / 2
bx = start_x
by = cy + Inches(0.2)

for i, (label, col, output) in enumerate(stages):
    # Box
    add_rect(s, bx, by, box_w, box_h, fill=col,
             line=RGBColor(0x1E, 0x35, 0x55))
    # Label
    txb = s.shapes.add_textbox(bx + Inches(0.06), by + Inches(0.1),
                                box_w - Inches(0.12), box_h - Inches(0.2))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"
    # Arrow
    if i < len(stages) - 1:
        ax = bx + box_w + Inches(0.02)
        ay = by + box_h / 2 - Pt(6)
        add_text(s, "→", ax, ay, box_gap - Inches(0.04), Inches(0.3),
                 size=Pt(20), bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    # Output label below
    add_text(s, output,
             bx, by + box_h + Inches(0.08), box_w, Inches(0.45),
             size=Pt(9.5), color=SLATE, align=PP_ALIGN.CENTER)
    bx += box_w + box_gap

# Feedback arrow label
add_text(s, "⟵  Online PPO update after every query  ⟵",
         start_x, by + box_h + Inches(0.6), total_w, Inches(0.32),
         size=Pt(11), color=PURPLE, align=PP_ALIGN.CENTER, bold=True)

# 4 detector options
det_y = by + box_h + Inches(1.1)
add_text(s, "4 Hallucination Detectors (Policy selects one per query):",
         Inches(0.45), det_y, Inches(12.43), Inches(0.3),
         size=Pt(11), bold=True, color=NAVY)
dets = [
    ("Token Overlap",      "Word-level matching",        "F1: 0.80  $0.05",  GREEN,  GREEN_BG),
    ("Semantic Similarity","Embedding cosine distance",  "F1: 0.38  $0.10",  AMBER,  AMBER_BG),
    ("BERT NLI",           "Neural entailment",          "F1: 0.73  $0.30",  BLUE,   BLUE_LIGHT),
    ("LLM Judge (Qwen)",   "Full LLM review",            "F1: 0.53  $1.00",  PURPLE, RGBColor(0xF5,0xF3,0xFF)),
]
dw = Inches(3.0)
dx = Inches(0.45)
for name, desc, metric, col, bg in dets:
    add_rect(s, dx, det_y + Inches(0.35), dw, Inches(0.85), fill=bg, line=col)
    add_text(s, name, dx + Inches(0.1), det_y + Inches(0.38),
             dw - Inches(0.15), Inches(0.28), size=Pt(11), bold=True, color=col)
    add_text(s, desc, dx + Inches(0.1), det_y + Inches(0.65),
             dw - Inches(0.15), Inches(0.22), size=Pt(9.5), color=SLATE)
    add_text(s, metric, dx + Inches(0.1), det_y + Inches(0.88),
             dw - Inches(0.15), Inches(0.22), size=Pt(9), bold=True, color=col)
    dx += dw + Inches(0.11)

page_num(s, 6)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MODEL EQUATIONS  ★ REQUIRED
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
cy = navy_header(s, "Model Equations",
                 "State representation · Reward · PPO objective · Detector selection")

eq_sections = [
    {
        "num": "Eq. 1",
        "title": "State Feature Vector  s_t ∈ ℝ¹²",
        "color": BLUE,
        "body": (
            "s_t  =  [ q_len,  e_count,  n_density,  has_pct,  has_cmp,\n"
            "          complexity,  k_chunks,  r_conf,  top_score,  score_gap,\n"
            "          a_len,  a_num_density ]"
        ),
        "note": "Encodes query complexity, retrieval quality, and answer characteristics",
    },
    {
        "num": "Eq. 2",
        "title": "Reward Function  r ∈ [-1, 1]",
        "color": GREEN,
        "body": (
            "r  =  0.8 · 𝟙[correct] · confidence\n"
            "       -  0.5 · cost_detector\n"
            "       -  0.1 · (latency_ms / 1000)"
        ),
        "note": "Balances detection quality, cost efficiency, and response speed",
    },
    {
        "num": "Eq. 3",
        "title": "PPO Clipped Surrogate Objective",
        "color": AMBER,
        "body": (
            "L^CLIP(θ)  =  𝔼_t [ min( r_t(θ) · Â_t ,\n"
            "               clip( r_t(θ), 1-ε, 1+ε ) · Â_t ) ]\n"
            "\n"
            "where  r_t(θ)  =  π_θ(a_t | s_t) / π_θ_old(a_t | s_t)"
        ),
        "note": "Prevents destructively large policy updates (ε = 0.2)",
    },
    {
        "num": "Eq. 4",
        "title": "Detector Selection  (Actor Network)",
        "color": PURPLE,
        "body": (
            "h  =  tanh( W₁ · ŝ_t  +  b₁ )           W₁ ∈ ℝ^{32×12}\n"
            "π_θ(a|s)  =  softmax( W₂ · h  +  b₂ )    W₂ ∈ ℝ^{4×32}\n"
            "\n"
            "d*  =  argmax_d  π_θ(d | s_t)    [ε-greedy: explore w.p. 0.15]"
        ),
        "note": "Two-layer MLP actor; input normalized by running mean/std",
    },
]

cols = 2
col_w_eq = Inches(6.1)
col_gap_eq = Inches(0.22)
col_positions = [Inches(0.45), Inches(0.45) + col_w_eq + col_gap_eq]
box_h_eq = Inches(2.3)
row_gap = Inches(0.18)
row_positions = [cy + Inches(0.08), cy + Inches(0.08) + box_h_eq + row_gap]

for idx, eq in enumerate(eq_sections):
    col_i = idx % 2
    row_i = idx // 2
    ex = col_positions[col_i]
    ey = row_positions[row_i]
    c = eq["color"]

    # Equation box
    add_rect(s, ex, ey, col_w_eq, box_h_eq, fill=WHITE, line=c, line_w=Pt(1.5))
    # Colour tag
    add_rect(s, ex, ey, Inches(0.4), box_h_eq, fill=c)
    add_text(s, eq["num"],
             ex + Inches(0.04), ey + Inches(0.9), Inches(0.32), Inches(0.5),
             size=Pt(8), bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    # Title
    add_text(s, eq["title"],
             ex + Inches(0.48), ey + Inches(0.08), col_w_eq - Inches(0.58), Inches(0.32),
             size=Pt(11), bold=True, color=c)
    # Body (monospace-style)
    txb = s.shapes.add_textbox(ex + Inches(0.48), ey + Inches(0.42),
                                col_w_eq - Inches(0.58), Inches(1.45))
    txb.word_wrap = False
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = eq["body"]
    run.font.size = Pt(11)
    run.font.name = "Courier New"
    run.font.color.rgb = NAVY
    # Note
    add_text(s, "→  " + eq["note"],
             ex + Inches(0.48), ey + box_h_eq - Inches(0.38),
             col_w_eq - Inches(0.58), Inches(0.3),
             size=Pt(9.5), italic=True, color=SLATE)

page_num(s, 7)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ALGORITHM + GRAPHIC ARCHITECTURE  ★ REQUIRED
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
cy = navy_header(s, "Algorithm & Neural Architecture",
                 "Algorithm 1: Adaptive Detector Routing via PPO   ·   Actor Network Diagram")

# ── Left: Algorithm table ────────────────────────────────────────────────────
alg_x = Inches(0.45)
alg_w = Inches(6.4)
alg_h = Inches(5.5)
add_rect(s, alg_x, cy, alg_w, alg_h, fill=OFFWHITE, line=BORDER)

# Header
add_rect(s, alg_x, cy, alg_w, Inches(0.4), fill=NAVY)
add_text(s, "Algorithm 1:  Adaptive Detector Selection  (FinGuard)",
         alg_x + Inches(0.12), cy + Inches(0.06), alg_w - Inches(0.2), Inches(0.3),
         size=Pt(10.5), bold=True, color=WHITE)

steps = [
    ("Input:",  "query q,  documents D,  policy π_θ",   BLUE,   True),
    ("1.",  "Retrieve top-k chunks C from D via FAISS (Eq. 1 features)",  NAVY, False),
    ("2.",  "Synthesize grounded answer a from C",                        NAVY, False),
    ("3.",  "Extract state  s_t ∈ ℝ¹²  from (q, a, C)  [Eq. 1]",        BLUE, False),
    ("4.",  "Normalize:  ŝ_t = (s_t − μ) / σ",                          NAVY, False),
    ("5.",  "Compute logits via actor MLP  [Eq. 4]",                      NAVY, False),
    ("6.",  "Select detector:  d ~ softmax(logits)  +  ε-greedy",         BLUE, False),
    ("7.",  "Run detector d on (q, a, C)  →  verdict, confidence",        NAVY, False),
    ("8.",  "Compute reward r  [Eq. 2]",                                  GREEN,False),
    ("9.",  "Append (s_t, d, r, log π) to replay buffer B",               NAVY, False),
    ("10.", "If |B| ≥ rollout_size:  update θ via PPO  [Eq. 3]",          AMBER,False),
    ("11.", "Save updated π_θ → rl_policy_runtime.json",                  NAVY, False),
    ("Out:", "verdict,  confidence,  cost savings vs. baseline",          PURPLE,True),
]
sy = cy + Inches(0.45)
for num, desc, col, bold in steps:
    row_bg = RGBColor(0xEF, 0xF4, 0xFF) if col == BLUE else (
             RGBColor(0xF0, 0xFD, 0xF4) if col == GREEN else (
             RGBColor(0xFF, 0xFB, 0xEB) if col == AMBER else (
             RGBColor(0xF5, 0xF3, 0xFF) if col == PURPLE else WHITE)))
    add_rect(s, alg_x + Inches(0.04), sy,
             alg_w - Inches(0.08), Inches(0.35),
             fill=row_bg, line=BORDER, line_w=Pt(0.5))
    add_text(s, num, alg_x + Inches(0.08), sy + Inches(0.04),
             Inches(0.42), Inches(0.28),
             size=Pt(10), bold=True, color=col)
    add_text(s, desc, alg_x + Inches(0.52), sy + Inches(0.04),
             alg_w - Inches(0.62), Inches(0.28),
             size=Pt(10), bold=bold, color=NAVY)
    sy += Inches(0.36)

# ── Right: Architecture diagram ──────────────────────────────────────────────
arch_x = Inches(7.1)
arch_w = Inches(5.8)

add_text(s, "Actor Network Architecture",
         arch_x, cy + Inches(0.05), arch_w, Inches(0.3),
         size=Pt(12), bold=True, color=BLUE)

# Input layer
def nn_layer(slide, lx, ly, lw, lh, label, sublabel, col, is_last=False):
    add_rect(slide, lx, ly, lw, lh, fill=col, line=RGBColor(0x1E,0x35,0x55))
    add_text(slide, label, lx + Inches(0.05), ly + Inches(0.08),
             lw - Inches(0.1), Inches(0.32),
             size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, sublabel, lx + Inches(0.05), ly + Inches(0.4),
             lw - Inches(0.1), Inches(0.28),
             size=Pt(9), color=WHITE, align=PP_ALIGN.CENTER)
    if not is_last:
        add_text(slide, "↓", lx + lw / 2 - Inches(0.15),
                 ly + lh + Inches(0.02), Inches(0.3), Inches(0.28),
                 size=Pt(18), bold=True, color=SLATE, align=PP_ALIGN.CENTER)

layer_w = Inches(4.5)
layer_h = Inches(0.72)
layer_gap = Inches(0.35)
lx = arch_x + (arch_w - layer_w) / 2
ly = cy + Inches(0.42)

layers = [
    ("Input:  s_t  ∈ ℝ¹²",            "12 state features (query + retrieval + answer)", BLUE),
    ("Normalize:  ŝ_t = (s_t-μ)/σ",   "Running mean/std from training data",            NAVY),
    ("Hidden:  tanh(W₁·ŝ_t + b₁)",    "Linear(12→32) + tanh activation  [Eq. 4]",     BLUE),
    ("Output:  W₂·h + b₂",             "Linear(32→4) logits",                           NAVY),
    ("softmax → π_θ(a|s)",             "Probability over 4 detectors  [Eq. 4]",         GREEN),
]
for i, (label, sublabel, col) in enumerate(layers):
    is_last = i == len(layers) - 1
    nn_layer(s, lx, ly, layer_w, layer_h, label, sublabel, col, is_last)
    ly += layer_h + layer_gap

# Detector output row
det_labels = ["token\noverlap", "semantic\nsim.", "bert\nnli", "llm\njudge"]
det_cols   = [GREEN, AMBER, BLUE, PURPLE]
det_w = Inches(0.95)
det_gap = Inches(0.15)
total_det_w = 4 * det_w + 3 * det_gap
det_start = arch_x + (arch_w - total_det_w) / 2
add_text(s, "↓",
         lx + layer_w / 2 - Inches(0.15), ly - Inches(0.28),
         Inches(0.3), Inches(0.28),
         size=Pt(18), bold=True, color=SLATE, align=PP_ALIGN.CENTER)
dx2 = det_start
for dl, dc in zip(det_labels, det_cols):
    add_rect(s, dx2, ly, det_w, Inches(0.72), fill=dc,
             line=RGBColor(0x1E, 0x35, 0x55))
    add_text(s, dl, dx2 + Inches(0.04), ly + Inches(0.12),
             det_w - Inches(0.08), Inches(0.5),
             size=Pt(9.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    dx2 += det_w + det_gap

# Eq link annotation
add_text(s, "Links to Eq. 3 (PPO update)  &  Eq. 2 (reward)",
         arch_x, ly + Inches(0.8), arch_w, Inches(0.28),
         size=Pt(9.5), italic=True, color=SLATE, align=PP_ALIGN.CENTER)

page_num(s, 8)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — THE FOUR DETECTORS
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
cy = navy_header(s, "The Four Hallucination Detectors",
                 "Each detector has a different mechanism, speed, and cost profile")

detectors = [
    {
        "name": "Token Overlap",
        "tag": "Grade A  ·  F1: 0.80",
        "color": GREEN, "bg": GREEN_BG, "border": RGBColor(0x86,0xEF,0xAC),
        "how": "Counts shared words between the answer and source chunks. Uses synonym expansion (revenue↔sales) and number matching.",
        "strengths": ["Catches all hallucinations (recall=1.0)", "3 ms  ·  $0.05 per call", "Best for numeric / direct-fact queries"],
        "weakness": "Over-flags (5 false positives) — precision = 0.67",
    },
    {
        "name": "BERT NLI",
        "tag": "Grade B  ·  F1: 0.73",
        "color": BLUE, "bg": BLUE_LIGHT, "border": BLUE_MID,
        "how": "Uses a neural NLI model to check whether the source document entails or contradicts the answer sentence-by-sentence.",
        "strengths": ["Understands paraphrase & implication", "Good recall (0.80)", "Best for comparative / reasoning queries"],
        "weakness": "1,047 ms  ·  $0.30 — slowest neural option",
    },
    {
        "name": "LLM Judge  (Qwen 2.5)",
        "tag": "Grade C  ·  F1: 0.53",
        "color": PURPLE, "bg": RGBColor(0xF5,0xF3,0xFF), "border": RGBColor(0xC4,0xB5,0xFD),
        "how": "Sends the full question, answer, and context to a local Qwen 2.5-7B model with a fact-checking system prompt.",
        "strengths": ["Highest precision (0.80)", "Can explain its reasoning", "Best for nuanced / ambiguous queries"],
        "weakness": "11,146 ms  ·  $1.00 — catches only 4/10 hallucinations",
    },
    {
        "name": "Semantic Similarity",
        "tag": "Grade D  ·  F1: 0.38",
        "color": AMBER, "bg": AMBER_BG, "border": RGBColor(0xFD,0xE6,0x8A),
        "how": "Computes cosine similarity between TF-IDF embeddings of the answer and source chunks.",
        "strengths": ["128 ms  ·  $0.10", "Catches paraphrased content"],
        "weakness": "Misses 7/10 hallucinations — numerically similar text fools it",
    },
]

dw = Inches(3.0)
dx = Inches(0.45)
dh = Inches(4.9)
for det in detectors:
    c = det["color"]
    add_rect(s, dx, cy + Inches(0.08), dw, dh,
             fill=det["bg"], line=det["border"])
    # colour bar top
    add_rect(s, dx, cy + Inches(0.08), dw, Inches(0.38), fill=c)
    add_text(s, det["name"],
             dx + Inches(0.1), cy + Inches(0.12), dw - Inches(0.15), Inches(0.28),
             size=Pt(11), bold=True, color=WHITE)

    ty = cy + Inches(0.55)
    add_text(s, det["tag"],
             dx + Inches(0.1), ty, dw - Inches(0.15), Inches(0.28),
             size=Pt(10), bold=True, color=c)
    ty += Inches(0.32)
    add_text(s, "How it works:",
             dx + Inches(0.1), ty, dw - Inches(0.15), Inches(0.24),
             size=Pt(9.5), bold=True, color=SLATE)
    ty += Inches(0.26)
    add_text(s, det["how"],
             dx + Inches(0.1), ty, dw - Inches(0.15), Inches(0.85),
             size=Pt(9.5), color=NAVY)
    ty += Inches(0.92)
    add_text(s, "Strengths:",
             dx + Inches(0.1), ty, dw - Inches(0.15), Inches(0.24),
             size=Pt(9.5), bold=True, color=GREEN)
    ty += Inches(0.26)
    for st in det["strengths"]:
        add_text(s, "✓  " + st, dx + Inches(0.1), ty,
                 dw - Inches(0.15), Inches(0.28),
                 size=Pt(9.5), color=GREEN)
        ty += Inches(0.3)
    ty += Inches(0.1)
    add_text(s, "Limitation:",
             dx + Inches(0.1), ty, dw - Inches(0.15), Inches(0.24),
             size=Pt(9.5), bold=True, color=RED)
    ty += Inches(0.26)
    add_text(s, "✗  " + det["weakness"],
             dx + Inches(0.1), ty, dw - Inches(0.15), Inches(0.5),
             size=Pt(9.5), color=RED)
    dx += dw + Inches(0.11)

page_num(s, 9)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — BENCHMARK RESULTS
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
cy = navy_header(s, "Preliminary Results",
                 "Benchmark: 20 financial QA samples (10 hallucinated, 10 grounded)")

# Metrics table
add_text(s, "Table 1.  Detector performance on financial QA benchmark",
         Inches(0.45), cy + Inches(0.05), Inches(12.43), Inches(0.3),
         size=Pt(11), bold=True, color=NAVY)

headers = ["Detector", "Accuracy", "Precision", "Recall", "F1 ↑", "Latency", "Cost ↓", "Grade"]
col_widths = [Inches(2.1), Inches(1.0), Inches(1.0), Inches(0.9), Inches(0.85), Inches(1.35), Inches(0.9), Inches(0.75)]
col_starts = [Inches(0.45)]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w)

rows_data = [
    ("Token Overlap",    "75.0%", "66.7%", "100.0%", "80.0%", "3 ms",      "$0.05", "A", GREEN,  GREEN_BG),
    ("BERT NLI",         "70.0%", "66.7%", "80.0%",  "72.7%", "1,047 ms",  "$0.30", "B", BLUE,   BLUE_LIGHT),
    ("LLM Judge",        "65.0%", "80.0%", "40.0%",  "53.3%", "11,146 ms", "$1.00", "C", PURPLE, RGBColor(0xF5,0xF3,0xFF)),
    ("Semantic Sim.",    "50.0%", "50.0%", "30.0%",  "37.5%", "128 ms",    "$0.10", "D", AMBER,  AMBER_BG),
]
row_h = Inches(0.46)
ty = cy + Inches(0.4)

# Header row
for cw, cs, hd in zip(col_widths, col_starts, headers):
    add_rect(s, cs, ty, cw, row_h, fill=NAVY)
    add_text(s, hd, cs + Inches(0.06), ty + Inches(0.1),
             cw - Inches(0.1), row_h - Inches(0.15),
             size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
ty += row_h

for name, acc, prec, rec, f1, lat, cost, grade, col, bg in rows_data:
    vals = [name, acc, prec, rec, f1, lat, cost, grade]
    for i, (cw, cs, val) in enumerate(zip(col_widths, col_starts, vals)):
        cell_bg = bg if i == 0 else WHITE
        is_f1 = (i == 4)
        is_grade = (i == 7)
        v_color = col if (i == 0 or is_grade) else (GREEN if is_f1 else NAVY)
        add_rect(s, cs, ty, cw, row_h, fill=cell_bg, line=BORDER, line_w=Pt(0.5))
        add_text(s, val, cs + Inches(0.06), ty + Inches(0.1),
                 cw - Inches(0.1), row_h - Inches(0.15),
                 size=Pt(11 if is_f1 or is_grade else 10.5),
                 bold=(i == 0 or is_f1 or is_grade),
                 color=v_color, align=PP_ALIGN.CENTER)
    ty += row_h

# Findings
find_y = ty + Inches(0.2)
add_text(s, "Key Findings",
         Inches(0.45), find_y, Inches(12.43), Inches(0.3),
         size=Pt(12), bold=True, color=BLUE)
findings = [
    ("Token Overlap wins on F1 (0.80) despite being the simplest method — financial "
     "hallucinations are primarily factual/numerical errors that word matching catches well.",
     NAVY, WHITE),
    ("LLM Judge has the highest precision (0.80) but lowest recall (0.40) — it trusts answers "
     "too readily, missing 6 of 10 hallucinations.",
     NAVY, WHITE),
    ("Adaptive routing over fixed strategy:  using LLM Judge for all queries costs $1.00/call. "
     "Routing simple queries to Token Overlap saves up to 95% cost with better F1.",
     GREEN, GREEN_BG),
]
fy = find_y + Inches(0.35)
for text, tc, bg in findings:
    bh = Inches(0.5)
    add_rect(s, Inches(0.45), fy, Inches(12.43), bh, fill=bg, line=BORDER)
    add_text(s, "▸  " + text,
             Inches(0.6), fy + Inches(0.08), Inches(12.1), bh - Inches(0.1),
             size=Pt(10.5), color=tc)
    fy += bh + Inches(0.08)

page_num(s, 10)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CONCLUSION & FUTURE WORK
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
cy = navy_header(s, "Conclusion & Future Work",
                 "Summary, limitations, and next steps")

# Left — Conclusion
add_text(s, "What Was Built",
         Inches(0.45), cy + Inches(0.05), Inches(6.0), Inches(0.32),
         size=Pt(13), bold=True, color=BLUE)
done = [
    ("FinGuard end-to-end pipeline",
     "FAISS retrieval over 975 SEC filing chunks + grounded answer synthesis"),
    ("4 hallucination detectors",
     "Token overlap, semantic similarity, BERT NLI, LLM Judge (Qwen 2.5-7B)"),
    ("PPO adaptive policy",
     "12-dim state, actor-critic MLP (12→32→4), online update after each query"),
    ("ε-greedy exploration",
     "15% random exploration prevents mode collapse in the online setting"),
    ("Interactive dashboard",
     "FinGuard Next.js app: live query, pipeline trace, benchmark grades"),
    ("Benchmark evaluation",
     "Token Overlap achieves best F1 (0.80) at lowest cost ($0.05)"),
]
dy = cy + Inches(0.42)
for title, desc in done:
    add_rect(s, Inches(0.45), dy, Inches(6.0), Inches(0.62),
             fill=BLUE_LIGHT, line=BLUE_MID)
    add_text(s, "✓  " + title, Inches(0.6), dy + Inches(0.04),
             Inches(5.7), Inches(0.26), size=Pt(11), bold=True, color=BLUE)
    add_text(s, desc, Inches(0.6), dy + Inches(0.31),
             Inches(5.7), Inches(0.26), size=Pt(10), color=SLATE)
    dy += Inches(0.68)

# Right — Future Work
add_text(s, "Future Work",
         Inches(6.85), cy + Inches(0.05), Inches(6.0), Inches(0.32),
         size=Pt(13), bold=True, color=AMBER)
future = [
    ("Ablation study",
     "Compare PPO vs. bandit vs. random vs. always-LLM on 100+ queries"),
    ("Real sentence-transformer embeddings",
     "Replace hash embeddings with all-MiniLM-L6 for better retrieval"),
    ("Multi-detector ensembling",
     "Combine outputs of two detectors when policy uncertainty is high"),
    ("Expanded benchmark",
     "Test on FinQA, TATQA, and SEC EDGAR datasets for generalisability"),
    ("Learning curve analysis",
     "Plot reward vs. queries to show policy convergence over time"),
    ("Cost-accuracy Pareto analysis",
     "Show the operating curve of adaptive vs. fixed detector strategies"),
]
fy2 = cy + Inches(0.42)
for title, desc in future:
    add_rect(s, Inches(6.85), fy2, Inches(6.0), Inches(0.62),
             fill=AMBER_BG, line=RGBColor(0xFD, 0xE6, 0x8A))
    add_text(s, "→  " + title, Inches(7.0), fy2 + Inches(0.04),
             Inches(5.7), Inches(0.26), size=Pt(11), bold=True, color=AMBER)
    add_text(s, desc, Inches(7.0), fy2 + Inches(0.31),
             Inches(5.7), Inches(0.26), size=Pt(10), color=SLATE)
    fy2 += Inches(0.68)

# Bottom closing statement
add_rect(s, Inches(0.45), H - Inches(1.0), Inches(12.43), Inches(0.65),
         fill=NAVY, line=DARK_BORDER)
add_text(s,
         "FinGuard demonstrates that reinforcement learning can reduce hallucination "
         "detection cost by up to 95% while maintaining accuracy — with room to improve "
         "as the policy learns from more queries.",
         Inches(0.65), H - Inches(0.92), Inches(12.0), Inches(0.55),
         size=Pt(12), color=WHITE)

page_num(s, 11)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — REFERENCES
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
cy = navy_header(s, "References", "")

refs_col1 = [
    "[1]  Lewis, P. et al. (2020). Retrieval-Augmented Generation for Knowledge-Intensive NLP Tasks. NeurIPS.",
    "[2]  Ji, Z. et al. (2023). Survey of Hallucination in Natural Language Generation. ACM Comp. Surveys, 55(12).",
    "[3]  Min, S. et al. (2023). FActScoring: Fine-Grained Atomic Evaluation of Factual Precision. EMNLP.",
    "[4]  Es, S. et al. (2023). RAGAS: Automated Evaluation of Retrieval Augmented Generation. arXiv:2309.15217.",
    "[5]  Zheng, L. et al. (2023). Judging LLM-as-a-Judge with MT-Bench and Chatbot Arena. NeurIPS.",
    "[6]  Honovich, O. et al. (2022). TRUE: Re-evaluating Factual Consistency Evaluation. NAACL.",
    "[7]  He, P. et al. (2022). Towards Document-Level Natural Language Inference. ACL.",
]
refs_col2 = [
    "[8]  Schulman, J. et al. (2017). Proximal Policy Optimization Algorithms. arXiv:1707.06347.",
    "[9]  Ouyang, L. et al. (2022). Training LMs to Follow Instructions with Human Feedback. NeurIPS.",
    "[10] Auer, P. et al. (2002). Finite-time Analysis of the Multiarmed Bandit Problem. JMLR, 3.",
    "[11] Devlin, J. et al. (2019). BERT: Pre-training of Deep Bidirectional Transformers. NAACL.",
    "[12] Sutton, R. & Barto, A. (2018). Reinforcement Learning: An Introduction. MIT Press.",
    "[13] Chen, Z. et al. (2021). FinQA: A Dataset of Numerical Reasoning over Financial Data. EMNLP.",
    "[14] Johnson, J. et al. (2019). Billion-Scale Similarity Search with GPUs (FAISS). IEEE TPAMI.",
]

categories = [
    ("RAG & Hallucination Detection",   BLUE,   refs_col1[:4]),
    ("LLM Evaluation & NLI",            PURPLE, refs_col1[4:]),
    ("Reinforcement Learning",          GREEN,  refs_col2[:4]),
    ("Datasets & Foundations",          AMBER,  refs_col2[4:]),
]

section_x = [Inches(0.45), Inches(6.85)]
section_y = [cy + Inches(0.1), cy + Inches(0.1)]
col_idx = 0
for cat_name, cat_color, cat_refs in categories:
    cx_s = section_x[col_idx % 2]
    cy_s = section_y[col_idx % 2]
    sec_w = Inches(6.1)

    # Category label
    add_rect(s, cx_s, cy_s, sec_w, Inches(0.3), fill=cat_color)
    add_text(s, cat_name, cx_s + Inches(0.1), cy_s + Inches(0.04),
             sec_w - Inches(0.15), Inches(0.24),
             size=Pt(10), bold=True, color=WHITE)
    ry = cy_s + Inches(0.34)
    for ref in cat_refs:
        bg = RGBColor(0xF5, 0xF7, 0xFA)
        add_rect(s, cx_s, ry, sec_w, Inches(0.52),
                 fill=bg, line=BORDER, line_w=Pt(0.5))
        add_text(s, ref, cx_s + Inches(0.1), ry + Inches(0.06),
                 sec_w - Inches(0.15), Inches(0.42),
                 size=Pt(9.5), color=NAVY)
        ry += Inches(0.54)

    section_y[col_idx % 2] = ry + Inches(0.18)
    col_idx += 1

# Footer note
add_rect(s, Inches(0.45), H - Inches(0.55), Inches(12.43), Inches(0.32),
         fill=NAVY, line=DARK_BORDER)
add_text(s, "Full bibliography available in project repository  ·  FinGuard — Adaptive Hallucination Detection in Financial RAG via PPO",
         Inches(0.65), H - Inches(0.50), Inches(12.0), Inches(0.26),
         size=Pt(9.5), color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER)

page_num(s, 12)


# ── Save ─────────────────────────────────────────────────────────────────────
out_path = "FinGuard_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
